from pptx import Presentation
import copy

from core_funcs import retrieve_ppt_s3

def load_template():
    """
    Get template from S3, load into Presentation
    """
    template_obj = retrieve_ppt_s3()
    prs = Presentation(template_obj)
    return prs


def _get_blank_slide_layout(pres):
    """
    Sourced from: https://github.com/scanny/python-pptx/issues/132
    """
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def duplicate_slide(pres, index: int):
    """
    Duplicate the slide with the given index in pres.
    Adds slide to the end of the presentation
    
    Sourced from: https://github.com/scanny/python-pptx/issues/132
    """
    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres)
    new_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    return new_slide


def delete_slide(presentation, index):
    """
    Deletes a slide at given index.
    
    Sources from: https://github.com/scanny/python-pptx/issues/67
    """
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[index])